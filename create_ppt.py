"""
IT人材派遣会社 事業計画 PPT 生成スクリプト
Using Claude AI Agent + Nano Banana Pro images
"""
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
IMG_DIR = os.path.join(os.path.dirname(__file__), "ppt_images")
OUTPUT = os.path.join(os.path.dirname(__file__), "IT人材派遣会社_事業計画.pptx")

# Colors
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_BLUE = RGBColor(0x1B, 0x4D, 0x89)
LIGHT_BLUE = RGBColor(0x41, 0x8F, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
GOLD = RGBColor(0xE8, 0xA8, 0x38)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def add_image(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        else:
            slide.shapes.add_picture(path, left, top)


# ========================================
# SLIDE 1: Cover
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)
add_image(slide, "cover.png", Inches(0), Inches(0), W, H)
# Dark overlay
overlay = add_shape(slide, Inches(0), Inches(0), W, H, DARK_BLUE)
overlay.fill.solid()
overlay.fill.fore_color.rgb = DARK_BLUE
# Title text
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
         "AI-Powered IT Staffing Company in Japan", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
         "IT\u4eba\u6750\u6d3e\u9063\u4e8b\u696d\u8a08\u753b\u66f8 \u2014 Claude Agent \u00d7 AI\u81ea\u52d5\u5316\u3067\u6b21\u4e16\u4ee3\u306e\u4eba\u6750\u30d3\u30b8\u30cd\u30b9\u3092",
         24, GOLD, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         "Focus: AI Development | UI/UX Design | Frontend | Backend | Java",
         20, LIGHT_BLUE, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "2026 Business Plan", 16, RGBColor(0x99, 0x99, 0x99), False, PP_ALIGN.CENTER)


# ========================================
# SLIDE 2: Market Opportunity
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
         "\u5e02\u5834\u6a5f\u4f1a  Market Opportunity", 32, WHITE, True)

add_image(slide, "market_growth.png", Inches(7.5), Inches(1.5), Inches(5.5), Inches(5.5))

# Key stats
stats = [
    ("\u00a59\u5146+", "IT\u6d3e\u9063\u5e02\u5834\u898f\u6a21"),
    ("79\u4e07\u4eba", "2030\u5e74 IT\u4eba\u6750\u4e0d\u8db3"),
    ("$115B", "2030\u5e74 IT\u30b5\u30fc\u30d3\u30b9\u4e88\u6e2c"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.8 + i * 1.7)
    add_shape(slide, Inches(0.5), y, Inches(6.5), Inches(1.4), LIGHT_GRAY)
    add_text(slide, Inches(0.8), y + Inches(0.1), Inches(3), Inches(0.8), num, 36, ACCENT_BLUE, True)
    add_text(slide, Inches(0.8), y + Inches(0.8), Inches(5.5), Inches(0.5), label, 18, DARK_TEXT)

add_bullet_text(slide, Inches(0.5), Inches(6.2), Inches(6.5), Inches(1),
                ["\u2022 METI\u4e88\u6e2c: 2025\u5e74\u306e\u5d16 \u2192 IT\u4eba\u6750\u5371\u6a5f\u304c\u6df1\u523b\u5316",
                 "\u2022 DX\u63a8\u9032\u3067AI/\u30af\u30e9\u30a6\u30c9\u4eba\u6750\u306e\u9700\u8981\u304c\u6025\u62e1\u5927"],
                14, DARK_TEXT)


# ========================================
# SLIDE 3: Billing Rates by Specialization
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
         "\u5358\u4fa1\u76f8\u5834\u3068\u53ce\u76ca\u30e2\u30c7\u30eb  Revenue Model", 32, WHITE, True)

# Table data
roles = [
    ("AI/ML Engineer", "\u00a5600K - \u00a51,000K+", "\u00a5315K", "\u00a5180K", "\u2605\u2605\u2605\u2605\u2605"),
    ("Java Developer", "\u00a5500K - \u00a5800K", "\u00a5245K", "\u00a5140K", "\u2605\u2605\u2605\u2605"),
    ("Backend (Python/Go)", "\u00a5450K - \u00a5700K", "\u00a5228K", "\u00a5130K", "\u2605\u2605\u2605\u2605"),
    ("Frontend (React/Vue)", "\u00a5400K - \u00a5750K", "\u00a5210K", "\u00a5115K", "\u2605\u2605\u2605\u2605"),
    ("UI/UX Designer", "\u00a5400K - \u00a5650K", "\u00a5184K", "\u00a5100K", "\u2605\u2605\u2605"),
]

# Header
headers = ["\u8077\u7a2e", "\u6708\u984d\u5358\u4fa1", "\u30de\u30fc\u30b8\u30f3(35%)", "\u7d14\u5229\u76ca/\u6708", "\u9700\u8981"]
col_widths = [Inches(2.5), Inches(2.5), Inches(2.2), Inches(2.2), Inches(2)]
x_start = Inches(0.8)

for j, (header, cw) in enumerate(zip(headers, col_widths)):
    x = x_start + sum(col_widths[:j])
    add_shape(slide, x, Inches(1.6), cw, Inches(0.6), ACCENT_BLUE)
    add_text(slide, x, Inches(1.65), cw, Inches(0.5), header, 16, WHITE, True, PP_ALIGN.CENTER)

for i, (role, rate, margin, profit, demand) in enumerate(roles):
    y = Inches(2.2 + i * 0.9)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    vals = [role, rate, margin, profit, demand]
    for j, (val, cw) in enumerate(zip(vals, col_widths)):
        x = x_start + sum(col_widths[:j])
        add_shape(slide, x, y, cw, Inches(0.8), bg_color)
        fc = GREEN if j == 3 else DARK_TEXT
        add_text(slide, x, y + Inches(0.15), cw, Inches(0.5), val, 15, fc, j == 3, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.5),
         "\u203b \u30de\u30fc\u30b8\u30f3\u7387 35% (SES\u30e2\u30c7\u30eb) | \u696d\u754c\u5e73\u5747 30.4% | AI\u81ea\u52d5\u5316\u3067\u30d0\u30c3\u30af\u30aa\u30d5\u30a3\u30b9\u30b3\u30b9\u30c8 50-60%\u524a\u6e1b",
         13, RGBColor(0x88, 0x88, 0x88), False)


# ========================================
# SLIDE 4: Claude AI Agent Capabilities
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Claude AI Agent \u3067\u4e8b\u696d\u3092\u81ea\u52d5\u5316", 32, WHITE, True)
add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.5),
         "AI-Powered Operations: 50-60% Back-Office Cost Reduction", 18, GOLD)

add_image(slide, "recruitment_ai.png", Inches(8), Inches(1.5), Inches(5), Inches(5.5))

agents = [
    ("\U0001f916 Candidate Screening Agent", "\u5c65\u6b74\u66f8\u89e3\u6790\u3001\u30b9\u30ad\u30eb\u30de\u30c3\u30c1\u30f3\u30b0\u3001\u30b7\u30e7\u30fc\u30c8\u30ea\u30b9\u30c8\u81ea\u52d5\u751f\u6210", "70%\u6642\u9593\u524a\u6e1b"),
    ("\U0001f91d Client Matching Agent", "\u30a8\u30f3\u30b8\u30cb\u30a2\u3068\u6848\u4ef6\u306e\u6700\u9069\u30de\u30c3\u30c1\u30f3\u30b0", "60%\u6642\u9593\u524a\u6e1b"),
    ("\u2696\ufe0f Compliance Agent", "\u6d3e\u9063\u671f\u9593\u76e3\u8996\u3001\u4fdd\u967a\u72b6\u6cc1\u3001\u30ad\u30e3\u30ea\u30a2\u5f62\u6210\u652f\u63f4\u8a08\u753b", "\u5b8c\u5168\u81ea\u52d5\u5316"),
    ("\U0001f4ca Reporting Agent", "\u9031\u6b21/\u6708\u6b21KPI\u30ec\u30dd\u30fc\u30c8\u3001\u7a3c\u50cd\u7387\u5206\u6790", "80%\u6642\u9593\u524a\u6e1b"),
    ("\U0001f4e7 Communication Agent", "\u30af\u30e9\u30a4\u30a2\u30f3\u30c8/\u5019\u88dc\u8005\u5bfe\u5fdc\u3001\u65e5\u82f1\u4e21\u5bfe\u5fdc", "40%\u6642\u9593\u524a\u6e1b"),
]

for i, (name, desc, saving) in enumerate(agents):
    y = Inches(1.8 + i * 1.05)
    add_shape(slide, Inches(0.5), y, Inches(7), Inches(0.9), ACCENT_BLUE)
    add_text(slide, Inches(0.7), y + Inches(0.05), Inches(4), Inches(0.4), name, 16, WHITE, True)
    add_text(slide, Inches(0.7), y + Inches(0.45), Inches(4.5), Inches(0.4), desc, 12, RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, Inches(5.5), y + Inches(0.2), Inches(2), Inches(0.5), saving, 15, GOLD, True)


# ========================================
# SLIDE 5: Services & Team
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
         "\u30b5\u30fc\u30d3\u30b9\u9818\u57df  Service Areas", 32, WHITE, True)

add_image(slide, "ai_team.png", Inches(0.3), Inches(1.5), Inches(6), Inches(4))
add_image(slide, "services.png", Inches(6.5), Inches(1.5), Inches(6.5), Inches(4))

# Service boxes
services = [
    ("AI/ML\u958b\u767a", "LLM\u5c0e\u5165, \u6a5f\u68b0\u5b66\u7fd2, \u30c7\u30fc\u30bf\u5206\u6790"),
    ("Frontend", "React, Vue, Next.js, TypeScript"),
    ("Backend", "Java, Python, Go, Node.js"),
    ("UI/UX", "\u30c7\u30b6\u30a4\u30f3\u30b7\u30b9\u30c6\u30e0, Figma"),
    ("SES/\u6d3e\u9063", "\u30d7\u30ed\u30b8\u30a7\u30af\u30c8\u5358\u4f4d\u306e\u30c1\u30fc\u30e0\u63d0\u4f9b"),
]

for i, (svc, desc) in enumerate(services):
    x = Inches(0.5 + (i % 3) * 4.2)
    y = Inches(5.8) if i < 3 else Inches(6.5)
    add_shape(slide, x, y, Inches(3.8), Inches(0.6), ACCENT_BLUE)
    add_text(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.5), Inches(0.5), svc, 14, WHITE, True)
    add_text(slide, x + Inches(1.6), y + Inches(0.05), Inches(2.1), Inches(0.5), desc, 12, RGBColor(0xDD, 0xDD, 0xDD))


# ========================================
# SLIDE 6: Legal & Setup
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
         "\u8a2d\u7acb\u8981\u4ef6\u3068\u30b3\u30b9\u30c8  Setup Requirements", 32, WHITE, True)

# Requirements
reqs = [
    ("\u4f1a\u793e\u5f62\u614b", "\u682a\u5f0f\u4f1a\u793e (KK)", "\u4fe1\u7528\u529b\u304c\u9ad8\u304f\u3001\u8cc7\u91d1\u8abf\u9054\u53ef\u80fd"),
    ("\u8cc7\u672c\u91d1", "\u00a520,000,000\u4ee5\u4e0a", "\u6d3e\u9063\u8a31\u53ef\u306e\u6700\u4f4e\u8981\u4ef6"),
    ("\u73fe\u9810\u91d1", "\u00a515,000,000\u4ee5\u4e0a", "\u8cc7\u672c\u91d1\u306b\u542b\u3080"),
    ("\u30aa\u30d5\u30a3\u30b9", "20\u33a1\u4ee5\u4e0a", "\u4e8b\u696d\u5c02\u7528\u30b9\u30da\u30fc\u30b9"),
    ("\u6d3e\u9063\u5143\u8cac\u4efb\u8005", "\u8b1b\u7fd2\u4fee\u4e86\u8005", "3\u5e74\u4ee5\u4e0a\u306e\u96c7\u7528\u7ba1\u7406\u7d4c\u9a13"),
    ("\u8a31\u53ef\u7533\u8acb", "\u00a5210,000", "\u5be9\u67fb\u671f\u9593 2-3\u30f6\u6708"),
]

for i, (item, value, note) in enumerate(reqs):
    y = Inches(1.5 + i * 0.9)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.5), y, Inches(5.8), Inches(0.8), bg)
    add_text(slide, Inches(0.7), y + Inches(0.15), Inches(1.8), Inches(0.5), item, 15, ACCENT_BLUE, True)
    add_text(slide, Inches(2.5), y + Inches(0.15), Inches(2), Inches(0.5), value, 15, DARK_TEXT, True)
    add_text(slide, Inches(4.5), y + Inches(0.15), Inches(1.8), Inches(0.5), note, 12, RGBColor(0x88, 0x88, 0x88))

# Cost breakdown
add_text(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
         "\u521d\u671f\u6295\u8cc7\u5185\u8a33", 22, ACCENT_BLUE, True)
costs = [
    ("\u4f1a\u793e\u8a2d\u7acb", "\u00a5225,000"),
    ("\u6d3e\u9063\u8a31\u53ef\u7533\u8acb", "\u00a5219,000"),
    ("\u8cc7\u672c\u91d1", "\u00a520,000,000"),
    ("\u30aa\u30d5\u30a3\u30b9\u6574\u5099", "\u00a52,200,000"),
    ("\u30c6\u30af\u30ce\u30ed\u30b8\u30fc", "\u00a51,600,000"),
    ("\u904b\u8ee2\u8cc7\u91d1(6\u30f6\u6708)", "\u00a55,000,000"),
]
for i, (item, cost) in enumerate(costs):
    y = Inches(2.2 + i * 0.65)
    add_text(slide, Inches(7.2), y, Inches(3), Inches(0.5), item, 14, DARK_TEXT)
    add_text(slide, Inches(10.5), y, Inches(2), Inches(0.5), cost, 14, DARK_TEXT, False, PP_ALIGN.RIGHT)

add_shape(slide, Inches(7), Inches(6.2), Inches(5.8), Inches(0.7), GOLD)
add_text(slide, Inches(7.2), Inches(6.3), Inches(3), Inches(0.5), "\u5408\u8a08", 18, DARK_BLUE, True)
add_text(slide, Inches(10), Inches(6.3), Inches(2.5), Inches(0.5), "\u00a529,250,000", 20, DARK_BLUE, True, PP_ALIGN.RIGHT)


# ========================================
# SLIDE 7: Growth Timeline
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "\u6210\u9577\u30ed\u30fc\u30c9\u30de\u30c3\u30d7  Growth Timeline", 32, WHITE, True)

phases = [
    ("Phase 1", "Month 1-6", "\u6e96\u5099\u30fb\u8a2d\u7acb\u30fb SES\u958b\u59cb",
     ["\u682a\u5f0f\u4f1a\u793e\u8a2d\u7acb", "\u6d3e\u9063\u8a31\u53ef\u7533\u8acb", "SES\u3067\u5148\u884c\u55b6\u696d", "\u30a8\u30f3\u30b8\u30cb\u30a2 3-5\u540d"]),
    ("Phase 2", "Month 6-18", "\u62e1\u5927\u30fb\u6d3e\u9063\u958b\u59cb",
     ["\u6d3e\u9063\u8a31\u53ef\u53d6\u5f97", "\u30a8\u30f3\u30b8\u30cb\u30a2 15-20\u540d", "\u640d\u76ca\u5206\u5c90\u70b9\u9054\u6210", "AI\u30a8\u30fc\u30b8\u30a7\u30f3\u30c8\u672c\u683c\u7a3c\u50cd"]),
    ("Phase 3", "Month 18-36", "\u5b89\u5b9a\u6210\u9577",
     ["\u30a8\u30f3\u30b8\u30cb\u30a2 25-50\u540d", "\u5927\u624b\u4f01\u696d\u30a2\u30ab\u30a6\u30f3\u30c8", "\u6708\u5229\u76ca \u00a5100\u4e07+", "\u7b2c2\u30aa\u30d5\u30a3\u30b9\u691c\u8a0e"]),
]

for i, (phase, period, title, items) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    # Phase header
    add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(0.8), GOLD)
    add_text(slide, x + Inches(0.1), Inches(1.55), Inches(1.5), Inches(0.35), phase, 18, DARK_BLUE, True)
    add_text(slide, x + Inches(0.1), Inches(1.9), Inches(3.5), Inches(0.35), period, 13, DARK_BLUE)
    # Content
    add_shape(slide, x, Inches(2.3), Inches(3.8), Inches(3.5), ACCENT_BLUE)
    add_text(slide, x + Inches(0.2), Inches(2.4), Inches(3.4), Inches(0.5), title, 17, WHITE, True)
    item_text = "\n".join([f"\u2713 {item}" for item in items])
    add_text(slide, x + Inches(0.2), Inches(3.0), Inches(3.4), Inches(2.5), item_text, 14, RGBColor(0xDD, 0xDD, 0xDD))

# Revenue projection
add_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         "\u53ce\u76ca\u4e88\u6e2c: 15\u540d\u00d7\u00a5650K\u5e73\u5747\u5358\u4fa1 = \u6708\u5546 \u00a59,750,000 | \u7d14\u5229\u76ca\u7387 8-12% | \u9ed2\u5b57\u5316\u76ee\u6a19: 18-24\u30f6\u6708",
         16, GOLD, False, PP_ALIGN.CENTER)


# ========================================
# SLIDE 8: Competitive Advantage
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(1.2), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
         "\u7af6\u4e89\u512a\u4f4d\u6027  Competitive Advantage", 32, WHITE, True)

# Traditional vs AI-powered
add_text(slide, Inches(1), Inches(1.5), Inches(5), Inches(0.5),
         "\u5f93\u6765\u306e\u6d3e\u9063\u4f1a\u793e (50\u540d\u898f\u6a21)", 20, RED, True)
trad = [
    "\u25cf \u63a1\u7528\u62c5\u5f53 5-8\u540d",
    "\u25cf \u55b6\u696d 2-3\u540d",
    "\u25cf \u7ba1\u7406/\u30b3\u30f3\u30d7\u30e9\u30a4\u30a2\u30f3\u30b9 2\u540d",
    "\u25cf \u7814\u4fee\u30b3\u30fc\u30c7\u30a3\u30cd\u30fc\u30bf\u30fc 1-2\u540d",
    "\u25cf \u30d0\u30c3\u30af\u30aa\u30d5\u30a3\u30b9\u5408\u8a08: 10-15\u540d",
]
add_bullet_text(slide, Inches(1), Inches(2.2), Inches(5), Inches(3), trad, 16, DARK_TEXT)

add_text(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
         "AI\u99c6\u52d5\u578b\u6d3e\u9063\u4f1a\u793e (50\u540d\u898f\u6a21)", 20, GREEN, True)
ai_co = [
    "\u25cf \u63a1\u7528\u62c5\u5f53 2-3\u540d + Claude Agent",
    "\u25cf \u55b6\u696d 1-2\u540d + Claude Agent",
    "\u25cf \u30b3\u30f3\u30d7\u30e9\u30a4\u30a2\u30f3\u30b9: Claude Agent\u81ea\u52d5\u5316",
    "\u25cf \u7814\u4fee: Claude\u304c\u30ad\u30e3\u30ea\u30a2\u8a08\u753b\u81ea\u52d5\u751f\u6210",
    "\u25cf \u30d0\u30c3\u30af\u30aa\u30d5\u30a3\u30b9\u5408\u8a08: 4-6\u540d",
]
add_bullet_text(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(3), ai_co, 16, DARK_TEXT)

# Savings highlight
add_shape(slide, Inches(3), Inches(5.5), Inches(7), Inches(1.2), GOLD)
add_text(slide, Inches(3.2), Inches(5.6), Inches(6.5), Inches(0.5),
         "\u30d0\u30c3\u30af\u30aa\u30d5\u30a3\u30b9\u30b3\u30b9\u30c8 50-60% \u524a\u6e1b", 24, DARK_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3.2), Inches(6.1), Inches(6.5), Inches(0.4),
         "\u2192 \u30a8\u30f3\u30b8\u30cb\u30a2\u306b\u9ad8\u3044\u9084\u5143\u7387\u3092\u63d0\u4f9b \u2192 \u512a\u79c0\u306a\u4eba\u6750\u3092\u78ba\u4fdd \u2192 \u7af6\u4e89\u512a\u4f4d",
         14, DARK_BLUE, False, PP_ALIGN.CENTER)


# ========================================
# SLIDE 9: Action Plan
# ========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_image(slide, "success.png", Inches(8.5), Inches(1.2), Inches(4.5), Inches(5.5))

add_text(slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
         "\u30a2\u30af\u30b7\u30e7\u30f3\u30d7\u30e9\u30f3  Next Steps", 32, WHITE, True)

steps = [
    ("Week 1-2", "\u6d3e\u9063\u5143\u8cac\u4efb\u8005\u8b1b\u7fd2\u53d7\u8b1b (\u00a59,000, 6\u6642\u9593)"),
    ("Week 2-4", "\u682a\u5f0f\u4f1a\u793e\u8a2d\u7acb (\u96fb\u5b50\u5b9a\u6b3e\u3067\u00a54\u4e07\u7bc0\u7d04)"),
    ("Week 4-6", "\u30aa\u30d5\u30a3\u30b9\u78ba\u4fdd (20\u33a1+, \u6771\u4eac\u30d3\u30b8\u30cd\u30b9\u8857)"),
    ("Week 6-8", "\u52b4\u50cd\u8005\u6d3e\u9063\u4e8b\u696d\u8a31\u53ef\u7533\u8acb"),
    ("Month 2-4", "SES\u3067\u5148\u884c\u55b6\u696d\u958b\u59cb (\u8a31\u53ef\u4e0d\u8981)"),
    ("Month 3-4", "Claude AI Agent \u30a4\u30f3\u30d5\u30e9\u69cb\u7bc9"),
    ("Month 4-6", "\u6d3e\u9063\u8a31\u53ef\u53d6\u5f97 \u2192 \u6d3e\u9063\u55b6\u696d\u958b\u59cb"),
    ("Month 6-12", "\u30a8\u30f3\u30b8\u30cb\u30a2 10-15\u540d\u898f\u6a21\u3078"),
]

for i, (when, what) in enumerate(steps):
    y = Inches(1.3 + i * 0.72)
    add_shape(slide, Inches(0.5), y, Inches(2), Inches(0.55), GOLD)
    add_text(slide, Inches(0.6), y + Inches(0.08), Inches(1.8), Inches(0.4), when, 14, DARK_BLUE, True)
    add_text(slide, Inches(2.7), y + Inches(0.08), Inches(5.5), Inches(0.4), what, 15, WHITE)

add_text(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4),
         "Powered by Claude AI Agent \u00d7 Nano Banana Pro | 2026", 13, RGBColor(0x88, 0x88, 0x88))


# ========================================
# Save
# ========================================
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
